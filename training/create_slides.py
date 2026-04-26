from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG        = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xCC, 0xD6, 0xE0)
DIM       = RGBColor(0x77, 0x77, 0x77)
GREEN     = RGBColor(0x2D, 0xC6, 0x53)
ORANGE    = RGBColor(0xFF, 0x8C, 0x00)
RED       = RGBColor(0xFF, 0x44, 0x44)
ROW_A     = RGBColor(0x11, 0x2B, 0x3E)
ROW_B     = RGBColor(0x0A, 0x20, 0x30)
HDR       = RGBColor(0x00, 0x60, 0x7A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def bg(slide):
    # Full-slide background rectangle (more reliable across viewers)
    rect(slide, 0, 0, 13.33, 7.5, BG)

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = color

def txt(slide, text, l, t, w, h, size=13, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def title(slide, text):
    rect(slide, 0, 0, 13.33, 0.04, ACCENT)
    rect(slide, 0, 0.04, 13.33, 0.9, RGBColor(0x0A, 0x20, 0x30))
    txt(slide, text, 0.4, 0.1, 12.5, 0.78, size=32, bold=True, color=ACCENT)
    rect(slide, 0, 7.46, 13.33, 0.04, ACCENT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Experiments
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
bg(s2)
title(s2, "Experiments")

# Pretraining line
rect(s2, 0.4, 1.1, 12.5, 0.58, ROW_A)
txt(s2, "Pretrained EchoJEPA (ViT-L) from scratch on iCardio  ·  77k patients  ·  4.1M images  ·  200 epochs",
    0.6, 1.2, 12.1, 0.38, size=13, color=GRAY)

# LVEF section header
txt(s2, "LVEF REGRESSION  —  attentive probe on frozen encoder  (41% of labeled data)",
    0.4, 1.85, 12.5, 0.3, size=10, bold=True, color=ACCENT)

# Table
cx = [0.4, 5.2, 8.8, 11.5]
cw = [4.7, 3.5, 2.6, 1.65]

# header
for x, w, h in zip(cx, cw, ["Dataset", "Ours  (iCardio pretrain)", "EchoJEPA paper", "Gap"]):
    rect(s2, x, 2.18, w, 0.44, HDR)
    txt(s2, h, x+0.1, 2.22, w-0.15, 0.34, size=12, bold=True)

rows = [
    ("EchoNet-Dynamic",                "6.39 MAE", "5.59 MAE", "+0.80"),
    ("EchoNet-Pediatric  (zero-shot)", "8.28 MAE", "7.97 MAE", "+0.31"),
    ("EchoNet-Pediatric  (fine-tuned)","6.05 MAE", "5.50 MAE", "+0.55"),
    ("CAMUS",                          "running…", "    —",    "    —"),
]
for r, row in enumerate(rows):
    for i, (x, w, cell) in enumerate(zip(cx, cw, row)):
        rect(s2, x, 2.62 + r*0.5, w, 0.5, ROW_A if r%2==0 else ROW_B)
        c = WHITE if i == 1 and "MAE" in cell else \
            ORANGE if i == 3 and "+" in cell else GRAY
        txt(s2, cell, x+0.1, 2.68 + r*0.5, w-0.15, 0.38, size=12, color=c)

txt(s2, "EchoJEPA trained on 300k patients (18M videos) — gap is partly explained by 4× more data",
    0.4, 4.75, 12.5, 0.32, size=10, italic=True, color=DIM)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Proposal & Updated Plan
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
bg(s3)
title(s3, "Proposal & Updated Plan")

# Research question
rect(s3, 0.4, 1.08, 12.5, 0.5, RGBColor(0x05, 0x22, 0x35))
txt(s3,
    "Does JEPA-style pretraining yield similar performance & robustness benefits on a different clinical dataset?",
    0.55, 1.13, 12.1, 0.4, size=12, italic=True, color=ACCENT)

# Table columns
cx = [0.4,  4.3,  8.0,  10.6]
cw = [3.8,  3.6,  2.5,  2.55]

t0 = 1.72
rh = 1.08

for x, w, h in zip(cx, cw, ["Research Question", "How", "Status", "Update"]):
    rect(s3, x, t0, w, 0.4, HDR)
    txt(s3, h, x+0.1, t0+0.07, w-0.15, 0.28, size=12, bold=True)

table_rows = [
    (
        "Does JEPA work on\na different dataset?",
        "Train from scratch on iCardio\n→ evaluate on benchmarks",
        (GREEN,  "✓  Done"),
        "Competitive with EchoJEPA\nusing ~25% of their data",
        GREEN,
    ),
    (
        "Does JEPA transfer\nacross datasets?",
        "Fine-tune EchoJEPA's\nreleased weights on iCardio",
        (RED,    "✗  Eliminated"),
        "EchoJEPA did not release\nmodel weights",
        DIM,
    ),
    (
        "Is JEPA better than\nno pretraining?",
        "Compare vs. supervised ViT\ntrained on iCardio",
        (ORANGE, "→  Pending"),
        "Need to train\nsupervised baseline",
        ORANGE,
    ),
    (
        "Is JEPA better than\nother SSL methods?",
        "Compare vs. MAE\n(reconstruction SSL) on iCardio",
        (ORANGE, "→  Pending"),
        "Need to train\nMAE baseline",
        ORANGE,
    ),
]

for r, (q, how, (sc, sl), note, nc) in enumerate(table_rows):
    ty  = t0 + 0.4 + r * rh
    rbg = ROW_A if r % 2 == 0 else ROW_B
    dim = (sc == RED)
    qc  = DIM if dim else WHITE
    gc  = DIM if dim else GRAY

    for x, w in zip(cx, cw):
        rect(s3, x, ty, w, rh, rbg)

    txt(s3, q,    cx[0]+0.1, ty+0.15, cw[0]-0.15, rh-0.2, size=12, bold=True, color=qc)
    txt(s3, how,  cx[1]+0.1, ty+0.12, cw[1]-0.15, rh-0.16, size=11, color=gc)
    txt(s3, sl,   cx[2]+0.12, ty+0.3,  cw[2]-0.2,  0.42,  size=12, bold=True, color=sc)
    txt(s3, note, cx[3]+0.1, ty+0.15, cw[3]-0.15, rh-0.2, size=11, italic=True, color=nc)

out = "/home/ahmedaly/iCardio/EchoJEPAv2/training/progress_slides.pptx"
prs.save(out)
print(f"Saved: {out}")
