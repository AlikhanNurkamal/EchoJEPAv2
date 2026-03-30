#!/usr/bin/env python3
"""Generate iCardio Dataset & EchoJEPA onboarding presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Color palette --
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)      # dark navy
BG_SECTION = RGBColor(0x16, 0x21, 0x3E)    # section header bg
ACCENT = RGBColor(0x00, 0x96, 0xC7)        # blue accent
ACCENT2 = RGBColor(0xE0, 0x4F, 0x5F)       # red accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
TABLE_HEADER_BG = RGBColor(0x00, 0x56, 0x7A)
TABLE_ROW_BG = RGBColor(0xF0, 0xF7, 0xFA)
TABLE_ALT_BG = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, bullets, font_size=16,
                     color=WHITE, spacing=Pt(6), bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        if bold_first and ":" in bullet:
            # We can't do partial bold easily, so just bold the whole line
            pass
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None, font_size=12):
    """rows_data: list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.4 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
                    paragraph.alignment = PP_ALIGN.CENTER

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG

    return table_shape


# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Accent bar top
add_shape_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

# Title
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "iCardio Echocardiography Dataset", font_size=44, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "& EchoJEPA Self-Supervised Pretraining", font_size=32, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
            "Onboarding Presentation  |  Project Overview & Available Tasks",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape_rect(slide, 0, Inches(7.0), W, Inches(0.08), ACCENT)
add_textbox(slide, Inches(1), Inches(7.1), Inches(11), Inches(0.4),
            "MBZUAI  |  iCardio  |  2026", font_size=12, color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 2: Outline
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Outline", font_size=36, color=WHITE, bold=True)

sections = [
    "1.  Dataset Overview  --  Sources, structure, scale",
    "2.  Dataset Statistics  --  Studies, DICOMs, views, splits",
    "3.  Available Tasks  --  30+ classification & regression tasks",
    "4.  Task Deep-Dive  --  By heart structure (LV, RV, LA, RA, Valves, ...)",
    "5.  EchoJEPA  --  Self-supervised video pretraining",
    "6.  Current Progress  --  Preprocessing, training, evaluations",
    "7.  Opportunities  --  Downstream tasks & research directions",
]
add_bullet_slide(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(5),
                 sections, font_size=22, color=LIGHT_GRAY, spacing=Pt(14))


# ============================================================================
# SLIDE 3: Dataset Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Dataset Overview", font_size=36, color=WHITE, bold=True)

add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT)

# Key numbers
stats = [
    ("~80,000", "Studies", "Unique patient echocardiogram sessions"),
    ("~4.1M", "DICOMs", "Individual video/image sequences"),
    ("~2.1M", "Videos", "Standard (multi-frame) DICOMs"),
    ("28", "View Types", "A4C, PLAX, PSAX, A2C, A3C, ..."),
    ("30+", "Tasks", "Classification & regression from reports"),
]

for i, (num, label, desc) in enumerate(stats):
    x = Inches(0.8 + i * 2.4)
    y = Inches(1.6)
    add_textbox(slide, x, y, Inches(2.2), Inches(0.6), num, font_size=36, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.2), Inches(2.2), Inches(0.4), label, font_size=18, color=WHITE,
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.6), Inches(2.2), Inches(0.6), desc, font_size=12, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

# Data sources section
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
            "Data Sources", font_size=24, color=WHITE, bold=True)

sources = [
    "JSON Files: 16 study_info JSON files containing all EHR data, reports, diagnoses, measurements",
    "Each study includes: conditions, characterizations, stratifications, doctor notes, conclusions",
    "Each DICOM linked to: study UUID, view type, frame count, pixel spacing",
    "Labels extracted via REGEX from clinical text reports + direct measurements",
    "Train/Val/Test splits provided at the study level (study_designation field)",
]
add_bullet_slide(slide, Inches(1.2), Inches(4.1), Inches(11), Inches(3),
                 sources, font_size=15, color=LIGHT_GRAY, spacing=Pt(8))


# ============================================================================
# SLIDE 4: Dataset Structure
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Data Structure & Fields", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT)

# Study fields
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "Study-Level Fields", font_size=22, color=ACCENT, bold=True)

study_fields = [
    "study_uuid, study_designation (TRAIN/VAL/TEST)",
    "Demographics: age_at_visit, height, weight, BMI",
    "Diagnoses: conditions, characterizations, stratifications",
    "Doctor Notes: per-structure text (LV, RV, LA, RA, valves, ...)",
    "Measurements: EF, LVIDd, LVIDs, LA dimensions",
    "Doppler: E/A ratio, LVOT velocities, TR velocity, gradients",
    "Conclusions: free-text clinical summary",
]
add_bullet_slide(slide, Inches(1.0), Inches(2.0), Inches(5.5), Inches(4),
                 study_fields, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

# DICOM fields
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
            "DICOM-Level Fields", font_size=22, color=ACCENT, bold=True)

dicom_fields = [
    "dicom_uuid - unique identifier",
    "study_uuid - links to parent study",
    "type: Standard (video), Color (Doppler), Single Frame",
    "view: 28 echo view types (A4C, PLAX, PSAX, ...)",
    "n_frames: number of frames in the video",
    "physical_delta_x/y: pixel spacing (mm)",
]
add_bullet_slide(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4),
                 dicom_fields, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

# Preprocessed format
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
            "Preprocessed Format (WebDataset Shards)", font_size=22, color=ACCENT, bold=True)

prep_info = [
    "Each sample: (T, 336, 336, 3) uint8 RGB numpy array + JSON metadata",
    "Pipeline: Resample spacing -> Remove UI overlays -> Extract fan ROI -> Normalize -> Pad to square -> Resize 336px -> Resample to 24fps",
    "Stored as .tar shards (~1000 samples each, ~5.8GB per shard) for streaming training",
]
add_bullet_slide(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(1.5),
                 prep_info, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))


# ============================================================================
# SLIDE 5: View Distribution
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Echo View Distribution (28 Views)", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT)

views_data = [
    ["View", "Count", "View", "Count"],
    ["A4C", "119,272", "RVIT", "61,466"],
    ["PLAX Standard", "99,782", "SUB IVC", "61,285"],
    ["PSAX Mitral", "90,156", "A4C Zoomed LV", "54,887"],
    ["A3C", "85,073", "PSAX Apex", "48,125"],
    ["Unclear Dark", "83,807", "A2C Zoomed Mitral", "47,791"],
    ["A2C", "81,311", "Suprasternal Notch", "39,645"],
    ["PSAX Papillary", "75,173", "PSAXA Zoomed Aorta", "39,358"],
    ["A5C", "72,911", "PLAX Mitral Cusps", "36,988"],
    ["PSAX Zoomed Out", "72,070", "PLAX Aortic Cusps", "20,876"],
    ["Subcostal Standard", "68,634", "A5C Zoomed Aorta", "16,552"],
    ["PLAX Pericardial", "62,407", "A4C Zoomed RV", "9,754"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), views_data,
          col_widths=[Inches(3.2), Inches(1.8), Inches(3.7), Inches(1.8)], font_size=13)


# ============================================================================
# SLIDE 6: Tasks Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Available Tasks Overview", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT)

tasks_overview = [
    ["Structure", "Classification Tasks", "Regression Tasks"],
    ["Left Ventricle", "Enlargement, Hypertrophy,\nSystolic Function", "EF, LVIDd, LVIDs"],
    ["Left Atrium", "LAVI (severity)", "LAVI (mL/m\u00b2)"],
    ["Right Ventricle", "Enlargement, Systolic Function", "--"],
    ["Right Atrium", "Enlargement", "--"],
    ["Aortic Valve", "Stenosis, Regurgitation,\nBicuspid, Replacement", "Root dimension, Peak/Mean\nvelocity & gradient"],
    ["Mitral Valve", "Regurgitation, Stenosis,\nCalcification, MitralClip", "--"],
    ["Tricuspid Valve", "Regurgitation", "--"],
    ["Hemodynamics", "PA pressure, RA pressure,\nDilated IVC", "RVSP, PA pressure"],
    ["Diastolic Function", "Dysfunction grade (I-IV)", "E/A ratio"],
    ["Other", "Pericardial Effusion,\nPacemaker, Heart Failure", "--"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), tasks_overview,
          col_widths=[Inches(2.5), Inches(4.5), Inches(4.5)], font_size=12)


# ============================================================================
# SLIDE 7: LV Tasks - EF (the flagship task)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT2)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Left Ventricle: Ejection Fraction", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT2)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "The Most Common Echo AI Task", font_size=22, color=ACCENT, bold=True)

ef_desc = [
    "Percentage of blood pumped out each beat (normal: 50-70%)",
    "Requires temporal understanding of cardiac cycle",
    "79,398 labeled studies (94.3% from measurements, 19.5% from text)",
    "Regression target: continuous EF value",
    "5-number summary: 3, 58, 60, 65, 90",
    "Distribution peaks at 60-65% (healthy population bias)",
    "Gold-standard benchmark for echo AI models",
]
add_bullet_slide(slide, Inches(1.0), Inches(2.1), Inches(5.5), Inches(4),
                 ef_desc, font_size=15, color=LIGHT_GRAY, spacing=Pt(8))

# EF histogram text representation
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
            "EF Distribution", font_size=22, color=ACCENT, bold=True)

ef_dist = [
    ["EF Range", "Count", "%"],
    ["< 30%", "672", "0.8%"],
    ["30-40%", "986", "1.2%"],
    ["40-50%", "1,794", "2.3%"],
    ["50-55%", "3,700", "4.7%"],
    ["55-60%", "17,347", "21.8%"],
    ["60-65%", "30,397", "38.3%"],
    ["65-70%", "18,174", "22.9%"],
    ["> 70%", "6,328", "8.0%"],
]
add_table(slide, Inches(7), Inches(2.1), Inches(5.5), ef_dist,
          col_widths=[Inches(1.8), Inches(1.8), Inches(1.9)], font_size=13)

# Classification version
add_textbox(slide, Inches(7), Inches(5.5), Inches(5.5), Inches(0.4),
            "Also available as classification:", font_size=14, color=LIGHT_GRAY)
add_textbox(slide, Inches(7), Inches(5.9), Inches(5.5), Inches(1),
            "Normal (79.5%) | Mildly Reduced (9.1%) | Moderately Reduced (5.5%)\nSeverely Reduced (3.6%) | Reduced (2.4%)  --  11,238 samples",
            font_size=13, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 8: LV Tasks - Enlargement & Hypertrophy
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT2)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Left Ventricle: Enlargement & Hypertrophy", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT2)

# LV Enlargement
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "LV Enlargement / Dilation", font_size=22, color=ACCENT, bold=True)

enl_text = [
    "Chamber size increases beyond normal",
    "Based on internal diameter/volume during diastole & systole",
    "77,543 total samples",
]
add_bullet_slide(slide, Inches(1.0), Inches(2.1), Inches(5.5), Inches(1.5),
                 enl_text, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

enl_data = [
    ["Class", "Count", "%"],
    ["Normal", "75,380", "97.2%"],
    ["Enlargement", "1,332", "1.7%"],
    ["Mild", "522", "0.7%"],
    ["Severe", "179", "0.2%"],
    ["Moderate", "130", "0.2%"],
]
add_table(slide, Inches(0.8), Inches(3.3), Inches(5.5), enl_data,
          col_widths=[Inches(2), Inches(1.5), Inches(2)], font_size=12)

# LV Hypertrophy
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
            "LV Hypertrophy (Wall Thickening)", font_size=22, color=ACCENT, bold=True)

hyp_text = [
    "Thickening of the muscular wall",
    "Response to high blood pressure or aortic stenosis",
    "58,562 total samples",
]
add_bullet_slide(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(1.5),
                 hyp_text, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

hyp_data = [
    ["Class", "Count", "%"],
    ["Normal", "46,951", "80.2%"],
    ["Mild", "7,793", "13.3%"],
    ["Borderline", "2,483", "4.2%"],
    ["Moderate", "941", "1.6%"],
    ["Severe", "394", "0.7%"],
]
add_table(slide, Inches(7), Inches(3.3), Inches(5.5), hyp_data,
          col_widths=[Inches(2), Inches(1.5), Inches(2)], font_size=12)

# LVID regression
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
            "LV Internal Diameter (Regression): LVIDd (78,843 samples, median 4.8cm) | LVIDs (78,677 samples, median 3.2cm)",
            font_size=15, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 9: Diastolic Function & E/A
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT2)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Diastolic Function & E/A Ratio", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT2)

# Diastolic dysfunction
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "Diastolic Dysfunction (Classification)", font_size=22, color=ACCENT, bold=True)

dd_desc = [
    "How well the heart relaxes and fills with blood",
    "Graded I through IV by severity",
    "15,455 labeled samples",
    "Grade I: Impaired relaxation (mildest)",
    "Grade II: Pseudonormal filling (moderate)",
    "Grade III: Reversible restrictive (severe)",
]
add_bullet_slide(slide, Inches(1.0), Inches(2.1), Inches(5.5), Inches(2.5),
                 dd_desc, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

dd_data = [
    ["Class", "Count", "%"],
    ["Diastolic Dysfunction", "7,109", "46.0%"],
    ["Grade I", "4,392", "28.4%"],
    ["Grade II", "3,254", "21.1%"],
    ["Normal", "635", "4.1%"],
    ["Grade III", "65", "0.4%"],
]
add_table(slide, Inches(0.8), Inches(4.5), Inches(5.5), dd_data,
          col_widths=[Inches(2.5), Inches(1.5), Inches(1.5)], font_size=12)

# E/A ratio
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
            "E/A Ratio (Regression)", font_size=22, color=ACCENT, bold=True)

ea_desc = [
    "Ratio of early (E) to late (A) diastolic filling velocity",
    "Key indicator of diastolic function",
    "Normal: E/A > 1 in young adults",
    "1,487 labeled samples",
    "Average value: 1.3",
    "Range: 0-6 (most between 0.5-2.0)",
    "Extracted from Doppler measurements",
]
add_bullet_slide(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(3),
                 ea_desc, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))


# ============================================================================
# SLIDE 10: Left Atrium & Right Heart
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT2)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Left Atrium, Right Ventricle & Right Atrium", font_size=34, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT2)

# LA
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3.6), Inches(0.4),
            "Left Atrium Volume Index", font_size=20, color=ACCENT, bold=True)
la_items = [
    "Classification: 47,155 samples",
    "  Normal 56.4% | Mild 26.3%",
    "  Moderate 7.8% | Severe 8.8%",
    "Regression: 4,585 samples",
    "  Median: 34 mL/m\u00b2",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(3),
                 la_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# RV
add_textbox(slide, Inches(5), Inches(1.5), Inches(3.6), Inches(0.4),
            "Right Ventricle", font_size=20, color=ACCENT, bold=True)
rv_items = [
    "Enlargement: 78,150 samples",
    "  Normal 99.2% | Mild 0.6%",
    "  (Very imbalanced)",
    "Systolic Function: 9,501 samples",
    "  Normal 98.8% | Reduced 0.9%",
]
add_bullet_slide(slide, Inches(5), Inches(2.0), Inches(3.6), Inches(3),
                 rv_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# RA
add_textbox(slide, Inches(9.2), Inches(1.5), Inches(3.6), Inches(0.4),
            "Right Atrium", font_size=20, color=ACCENT, bold=True)
ra_items = [
    "Enlargement: 71,742 samples",
    "  Normal 90.1% | Mild 3.8%",
    "  Moderate 2.7% | Severe 1.6%",
    "  Enlarged 1.7%",
    "Better balanced than RV tasks",
]
add_bullet_slide(slide, Inches(9.2), Inches(2.0), Inches(3.6), Inches(3),
                 ra_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# Key challenge
add_shape_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.5), RGBColor(0x20, 0x30, 0x50))
add_textbox(slide, Inches(1.2), Inches(5.1), Inches(11), Inches(0.4),
            "Key Challenge: Class Imbalance", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.8),
            "Most tasks are heavily skewed toward 'Normal'. Strategies: weighted loss, focal loss, oversampling, "
            "few-shot learning, or collapse to binary (normal vs abnormal).",
            font_size=15, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 11: Aortic Valve Tasks
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT2)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Aortic Valve & Aorta Tasks", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT2)

av_tasks = [
    ["Task", "Type", "Samples", "Key Classes / Range"],
    ["Aortic Stenosis", "Classification", "62,888", "Normal 95.9% | Mild 2.0% | Moderate 1.6% | Severe 0.6%"],
    ["AS (iCardio CSV)", "Classification", "10,665", "Normal 79.4% | Mild 8.6% | Moderate 8.5% | Severe 3.5%"],
    ["Aortic Regurgitation", "Classification", "51,254", "Normal 80.9% | Trace 8.8% | Mild 8.1% | Mod 1.8%"],
    ["Bicuspid Aortic Valve", "Classification", "57,473", "Trileaflet 98.6% | Bicuspid 1.3%"],
    ["Aortic Root Dilation", "Classification", "76,334", "Normal 97.3% | Borderline 1.0% | Dilated 0.9%"],
    ["Aortic Root Dimension", "Regression", "2,379", "Range: 1.0-5.4 cm | Median: 4.0 cm"],
    ["Peak Velocity", "Regression", "710", "Range: 0.07-6.0 m/s | Median: 2.56 m/s"],
    ["Peak Gradient", "Regression", "927", "Range: 0.9-200 mmHg | Median: 28 mmHg"],
    ["Mean Gradient", "Regression", "1,887", "Range: 0.4-125 mmHg | Median: 20 mmHg"],
    ["AV Replacement", "Classification", "79,584", "No replacement 99.4% | Surgery 0.6%"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), av_tasks,
          col_widths=[Inches(2.2), Inches(1.5), Inches(1.3), Inches(7.3)], font_size=11)


# ============================================================================
# SLIDE 12: Mitral & Tricuspid Valve Tasks
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT2)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Mitral & Tricuspid Valve Tasks", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT2)

valve_tasks = [
    ["Task", "Type", "Samples", "Distribution"],
    ["Mitral Regurgitation", "Classification", "70,722", "Normal 44.4% | Mild 27.9% | Trace 23.9% | Mod 3.3% | Severe 0.5%"],
    ["Mitral Stenosis", "Classification", "79,584", "No stenosis 99.8% | Mild 0.1% | Present 0.1%"],
    ["Mitral Calcification", "Classification", "79,584", "No calc 95.4% | Present 2.2% | Mild 2.2% | Mod 0.3%"],
    ["MitralClip", "Classification", "79,584", "No intervention 99.97% | Present 25 samples (!)"],
    ["Tricuspid Regurgitation", "Classification", "72,530", "Normal 38.0% | Trace 30.5% | Mild 28.4% | Mod 2.8% | Severe 0.3%"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), valve_tasks,
          col_widths=[Inches(2.2), Inches(1.5), Inches(1.3), Inches(7.3)], font_size=12)

# Highlight
add_shape_rect(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.2), RGBColor(0x20, 0x30, 0x50))
add_textbox(slide, Inches(1.2), Inches(4.6), Inches(11), Inches(0.4),
            "Notable Observations", font_size=20, color=ACCENT, bold=True)
notable = [
    "Mitral Regurgitation has the most balanced distribution among valve tasks",
    "Tricuspid Regurgitation is also well-balanced (38/30/28% for normal/trace/mild)",
    "Mitral Stenosis and MitralClip are extremely rare -> few-shot or binary detection",
    "Multi-label complexity: many studies have regurgitation across multiple valves simultaneously",
]
add_bullet_slide(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(1.5),
                 notable, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))


# ============================================================================
# SLIDE 13: Hemodynamics & Other Tasks
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT2)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Hemodynamics, Pericardium & Devices", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT2)

hemo_tasks = [
    ["Task", "Type", "Samples", "Distribution / Range"],
    ["PA Pressure", "Classification", "53,666", "Normal 81.3% | Unable 12.9% | Mild 2.1% | Borderline 2.0%"],
    ["PA Pressure", "Regression", "5,799", "Range: 10-125 mmHg | Peak at 20-25 mmHg"],
    ["RVSP", "Regression", "885", "Range: 25-125 mmHg | Median ~48 mmHg"],
    ["RA Pressure", "Classification", "759", "Normal 90.5% | Elevated 3.3% | Sig. Elevated 6.2%"],
    ["Dilated IVC", "Classification", "9,128", "Normal 97.4% | Dilated 2.6%"],
    ["Pericardial Effusion", "Classification", "77,646", "Normal 99.4% | Tamponade 0.3% | Trace 0.2%"],
    ["Pacemaker", "Classification", "79,584", "No pacemaker 97.2% | Present 2.8%"],
    ["Heart Failure", "Classification", "10,117", "Normal 87.2% | Heart Failure 12.8%"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), hemo_tasks,
          col_widths=[Inches(2.2), Inches(1.5), Inches(1.3), Inches(7.3)], font_size=12)

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1),
            "Heart Failure is a multi-faceted task: includes HF with reduced EF (HFrEF), "
            "HF with preserved EF (HFpEF), and unspecified HF. Requires integration of "
            "multiple signals (EF, diastolic function, chamber sizes, filling pressures).",
            font_size=15, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 14: EchoJEPA Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
            "EchoJEPA", font_size=44, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(1.7), Inches(11), Inches(0.6),
            "Joint Embedding Predictive Architecture for Echocardiogram Video Understanding",
            font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.005), ACCENT)

# Two columns
add_textbox(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(0.4),
            "What is JEPA?", font_size=24, color=ACCENT, bold=True)

jepa_desc = [
    "Self-supervised pretraining (no labels needed)",
    "Learns by predicting masked video regions",
    "Encoder sees unmasked patches (context)",
    "Predictor predicts target encoder's representations",
    "Target encoder updated via EMA (momentum)",
    "Loss: L1 distance in representation space",
    "NOT pixel reconstruction (unlike MAE)",
    "Learns semantic features, not texture",
]
add_bullet_slide(slide, Inches(1.0), Inches(3.3), Inches(5.5), Inches(3.5),
                 jepa_desc, font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

add_textbox(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(0.4),
            "Why for Echo?", font_size=24, color=ACCENT, bold=True)

why_echo = [
    "Echocardiograms are inherently temporal (cardiac cycles)",
    "JEPA captures spatiotemporal structure",
    "Pre-train on millions of unlabeled echo videos",
    "Fine-tune on small labeled sets for downstream tasks",
    "Outperforms supervised baselines on EF, view cls, etc.",
    "Multi-view understanding (A4C, PLAX, PSAX, ...)",
    "Transfers across datasets (UHN, EchoNet, iCardio)",
]
add_bullet_slide(slide, Inches(7.2), Inches(3.3), Inches(5.5), Inches(3.5),
                 why_echo, font_size=15, color=LIGHT_GRAY, spacing=Pt(6))


# ============================================================================
# SLIDE 15: EchoJEPA Architecture
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "EchoJEPA: Architecture & Training", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT)

# Architecture
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
            "Architecture", font_size=22, color=ACCENT, bold=True)

arch = [
    "Encoder: ViT-Large/16 (304M params)",
    "Predictor: 12-layer transformer (22M params)",
    "Input: (B, C, 16, 336, 336) video clips",
    "Patch size: 16x16, Tubelet size: 2 (temporal)",
    "RoPE positional embeddings",
    "BFloat16 mixed precision",
]
add_bullet_slide(slide, Inches(1.0), Inches(2.0), Inches(5.5), Inches(3),
                 arch, font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# Masking
add_textbox(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.4),
            "Masking Strategy", font_size=22, color=ACCENT, bold=True)

masking = [
    "Context: 8 small spatiotemporal blocks (15% spatial)",
    "Target: 2 large blocks (70% spatial, full temporal)",
    "Encoder sees context, predictor predicts targets",
    "Masking ratio forces learning of global structure",
]
add_bullet_slide(slide, Inches(1.0), Inches(4.8), Inches(5.5), Inches(2),
                 masking, font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# Training
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Training Recipe", font_size=22, color=ACCENT, bold=True)

training = [
    "Phase 1 - Pretrain (200 epochs, 16 frames):",
    "  Constant LR after warmup (1.64e-5)",
    "  Batch size 24, 1 GPU (RTX A6000 48GB)",
    "  ~4.5s/iter, ~75 min/epoch",
    "",
    "Phase 2 - Cooldown (planned):",
    "  64 frames per clip (longer temporal context)",
    "  LR decay from peak to 1e-6",
    "  Refines temporal understanding",
    "",
    "Evaluation via frozen encoder:",
    "  Linear/Attentive probes on downstream tasks",
    "  No encoder fine-tuning needed",
]
add_bullet_slide(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4.5),
                 training, font_size=14, color=LIGHT_GRAY, spacing=Pt(4))


# ============================================================================
# SLIDE 16: Current Progress
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)
add_shape_rect(slide, 0, 0, Inches(0.15), H, ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Current Progress", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT)

# Preprocessing
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
            "Preprocessing", font_size=22, color=ACCENT, bold=True)

preproc = [
    "1.68M / 4.08M DICOMs processed (41%)",
    "2,288 shards created (~5.8GB each)",
    "Fan-ROI extraction + spacing normalization",
    "Stored as WebDataset .tar for streaming",
    "Resume capability via progress.csv",
]
add_bullet_slide(slide, Inches(1.0), Inches(2.0), Inches(5.5), Inches(2.5),
                 preproc, font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# Training
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "EchoJEPA Training", font_size=22, color=ACCENT, bold=True)

train_prog = [
    "Epoch ~93/200 (46.5% complete)",
    "Loss: 0.857 -> 0.442 (L1 in repr space)",
    "LR: reached peak 1.64e-5 (constant now)",
    "GPU memory: 17.1 GB / 48 GB",
    "Speed: ~4.5s/iter, ETA ~4-5 more days",
    "Wandb logging: all metrics tracked",
]
add_bullet_slide(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(2.5),
                 train_prog, font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# Eval results
add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
            "Preliminary Evaluation (Epoch 93-105)", font_size=22, color=ACCENT, bold=True)

eval_data = [
    ["Metric", "Value", "Notes"],
    ["Study-level Retrieval R@1", "62.2%", "Same-patient retrieval"],
    ["Study-level Retrieval R@5", "83.4%", ""],
    ["View Retrieval R@1", "54.8%", "Same echo view type"],
    ["View Classification Top-1", "15.4%", "24 classes, linear probe (3.7x random)"],
    ["View Classification Top-5", "51.7%", "Correct view often in top 5"],
    ["Feature Effective Rank", "91 / 1024", "Room for improvement (cooldown phase)"],
]
add_table(slide, Inches(0.8), Inches(5.3), Inches(11.5), eval_data,
          col_widths=[Inches(3.5), Inches(2), Inches(6)], font_size=12)


# ============================================================================
# SLIDE 17: Research Opportunities
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Research Opportunities", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.005), ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
            "Downstream Tasks", font_size=22, color=ACCENT, bold=True)

downstream = [
    "EF regression (benchmark against EchoNet/EchoPrime)",
    "Multi-task learning: predict EF + view + valve disease",
    "View classification (28 classes)",
    "Aortic stenosis severity classification",
    "Heart failure detection (HFrEF vs HFpEF)",
    "Rare condition detection (bicuspid valve, pacemaker)",
    "Report generation from video features",
]
add_bullet_slide(slide, Inches(1.0), Inches(2.0), Inches(5.5), Inches(3.5),
                 downstream, font_size=15, color=LIGHT_GRAY, spacing=Pt(8))

add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Methodology Research", font_size=22, color=ACCENT, bold=True)

methodology = [
    "Self-supervised pre-training strategies",
    "Handling extreme class imbalance",
    "Multi-view fusion (combine A4C + PLAX + PSAX)",
    "Temporal modeling for cardiac cycle understanding",
    "Few-shot learning for rare conditions",
    "Cross-dataset generalization (iCardio -> EchoNet)",
    "Retrieval-augmented diagnosis",
]
add_bullet_slide(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(3.5),
                 methodology, font_size=15, color=LIGHT_GRAY, spacing=Pt(8))

# Bottom callout
add_shape_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), RGBColor(0x10, 0x30, 0x45))
add_textbox(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.4),
            "Dataset Scale Advantage", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(1.2), Inches(6.3), Inches(11), Inches(0.6),
            "With ~80K studies and 4.1M DICOMs, this is one of the largest labeled echo datasets available. "
            "EchoNet-Dynamic has ~10K, CAMUS has ~500. The scale enables both strong pretraining "
            "and robust supervised baselines.",
            font_size=16, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 18: Thank You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
            "Questions?", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
            "Key Resources", font_size=24, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

resources = [
    "Code: /home/ahmedaly/iCardio/EchoJEPAv2/",
    "Preprocessing: preprocessing_alikhan/",
    "Training configs: training/pretrain_icardio_336px_16f.yaml",
    "Shards: /hdd2/ and /hdd1/ preprocessed_by_alikhan_for_echojepa/",
    "Checkpoints: /home/ahmedaly/iCardio/checkpoints/pretrain/",
    "Wandb: echojepa-pretrain project",
]
add_bullet_slide(slide, Inches(3), Inches(4.7), Inches(7), Inches(2.5),
                 resources, font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

add_shape_rect(slide, 0, Inches(7.0), W, Inches(0.08), ACCENT)

# ============================================================================
# Save
# ============================================================================
out_path = "/home/ahmedaly/iCardio/EchoJEPAv2/training/iCardio_EchoJEPA_Onboarding.pptx"
prs.save(out_path)
print(f"Presentation saved to: {out_path}")
print(f"Total slides: {len(prs.slides)}")
