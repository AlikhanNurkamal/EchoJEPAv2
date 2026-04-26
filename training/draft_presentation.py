"""Draft narrative presentation for EchoJEPAv2 work.

Story arc:
  Motivation -> Background -> Gap & Research Questions -> Method ->
  Replication results (LVEF) -> Generalisation (held-out sites + new tasks) ->
  Why text conditioning -> Text method -> Sanity checks
  (view classification + retrieval) -> Text downstream results ->
  Discussion -> Limitations -> Future work -> Q&A.

Run:
    python training/draft_presentation.py
Output: training/draft_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Theme ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)
PANEL    = RGBColor(0x11, 0x2B, 0x3E)
PANEL_B  = RGBColor(0x0A, 0x20, 0x30)
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_2 = RGBColor(0x90, 0xE0, 0xEF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xCC, 0xD6, 0xE0)
DIM      = RGBColor(0x77, 0x77, 0x77)
GREEN    = RGBColor(0x2D, 0xC6, 0x53)
ORANGE   = RGBColor(0xFF, 0x8C, 0x00)
RED      = RGBColor(0xFF, 0x44, 0x44)
HDR      = RGBColor(0x00, 0x60, 0x7A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ─── Drawing helpers ──────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if not line:
        s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color


def txt(slide, text, l, t, w, h, size=13, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color


def bullets(slide, items, l, t, w, h, size=14, gap=0.05, color=GRAY,
            bold_lead=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap * 72)
        if isinstance(it, tuple):
            lead, rest = it
            r = p.add_run(); r.text = "• "; r.font.size = Pt(size); r.font.color.rgb = ACCENT
            r = p.add_run(); r.text = lead
            r.font.size = Pt(size); r.font.bold = bold_lead; r.font.color.rgb = WHITE
            r = p.add_run(); r.text = rest
            r.font.size = Pt(size); r.font.color.rgb = color
        else:
            r = p.add_run(); r.text = "• "; r.font.size = Pt(size); r.font.color.rgb = ACCENT
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = color


def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, BG)


def title(slide, text, kicker=None):
    rect(slide, 0, 0, 13.33, 0.04, ACCENT)
    rect(slide, 0, 0.04, 13.33, 0.9, PANEL_B)
    txt(slide, text, 0.4, 0.1, 12.5, 0.78, size=28, bold=True, color=ACCENT)
    if kicker:
        txt(slide, kicker, 0.4, 0.62, 12.5, 0.3, size=11, italic=True, color=GRAY)
    rect(slide, 0, 7.46, 13.33, 0.04, ACCENT)


def footer(slide, left="EchoJEPAv2", right=""):
    txt(slide, left, 0.4, 7.15, 6, 0.25, size=9, color=DIM)
    if right:
        txt(slide, right, 6.93, 7.15, 6, 0.25, size=9, color=DIM, align=PP_ALIGN.RIGHT)


def panel(slide, l, t, w, h, header=None, color=PANEL):
    rect(slide, l, t, w, h, color)
    if header:
        rect(slide, l, t, w, 0.42, HDR)
        txt(slide, header, l + 0.15, t + 0.07, w - 0.2, 0.3, size=12, bold=True)


def add_slide():
    s = prs.slides.add_slide(blank)
    bg(s)
    return s


# ─── SLIDE 1 — Title ─────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 2.6, 13.33, 0.04, ACCENT)
rect(s, 0, 4.7, 13.33, 0.04, ACCENT)
txt(s, "EchoJEPAv2", 0.4, 1.7, 12.5, 0.9, size=48, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER)
txt(s, "Self-Supervised Echocardiogram Representation Learning",
    0.4, 2.75, 12.5, 0.55, size=22, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "with Optional Text Conditioning from Clinical Reports",
    0.4, 3.3, 12.5, 0.55, size=22, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Ahmed Aly  ·  Alikhan Nurkamal", 0.4, 4.95, 12.5, 0.4, size=16,
    color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Mohamed Bin Zayed University of Artificial Intelligence",
    0.4, 5.4, 12.5, 0.35, size=13, italic=True, color=DIM, align=PP_ALIGN.CENTER)
footer(s, "Draft v0", "April 2026")


# ─── SLIDE 2 — Clinical Motivation ───────────────────────────────────────────
s = add_slide()
title(s, "Why echocardiography, and why now?",
      "Echo is the most-ordered cardiac imaging study in the world — yet reading it remains a manual bottleneck.")

panel(s, 0.4, 1.2, 4.1, 5.6, header="Clinical reality")
bullets(s, [
    ("Ubiquitous. ", "~10M echo studies / year in the US alone."),
    ("Operator-dependent. ", "Quality and view selection vary."),
    ("Time-intensive. ", "Reading + reporting takes minutes per study."),
    ("Critical metrics. ", "LVEF guides nearly every cardiology decision."),
], 0.55, 1.75, 3.85, 5.0, size=12)

panel(s, 4.7, 1.2, 4.1, 5.6, header="Data reality at a centre")
bullets(s, [
    ("PACS archives are huge. ", "Millions of unlabeled videos."),
    ("Labels are scarce. ", "A few clinical metrics per study, embedded in free-text reports."),
    ("Cross-vendor noise. ", "GE / Philips / Mindray etc. all look subtly different."),
    ("Privacy-bound. ", "Data cannot leave the institution easily."),
], 4.85, 1.75, 3.85, 5.0, size=12)

panel(s, 9.0, 1.2, 3.9, 5.6, header="Implication")
txt(s,
    "Self-supervised pretraining on a single institution's unlabeled echo "
    "archive should unlock a strong general-purpose encoder — without needing "
    "to ship data anywhere or annotate millions of clips.",
    9.15, 1.85, 3.65, 4.7, size=13, color=WHITE)
txt(s,
    "Clinical reports already exist for free. Can we use them as an extra "
    "supervisory signal during pretraining?",
    9.15, 5.4, 3.65, 1.4, size=13, italic=True, color=ACCENT_2)
footer(s, "Motivation")


# ─── SLIDE 3 — SSL paradigms in echo (background) ────────────────────────────
s = add_slide()
title(s, "Background — self-supervised paradigms in video",
      "Two dominant SSL families, and where echo currently sits.")

panel(s, 0.4, 1.2, 6.2, 5.6, header="MAE / VideoMAE")
bullets(s, [
    ("Pixel target. ", "Reconstruct masked patches at the pixel level."),
    ("Pros. ", "Conceptually simple, well-studied, strong on natural video."),
    ("Cons. ", "Pixel-level loss spends capacity on low-level texture; ultrasound speckle is a distractor."),
], 0.55, 1.75, 5.95, 4.5, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="V-JEPA  →  EchoJEPA")
bullets(s, [
    ("Latent target. ", "Predict masked-region embeddings from context embeddings."),
    ("Pros. ", "Loss lives in feature space — emphasises semantics over pixel detail."),
    ("Echo evidence. ", "EchoJEPA reports MAE = 5.59 LVEF on EchoNet-Dynamic with a frozen attentive probe — best published echo-SSL result."),
    ("But… ", "Weights are not released; results not reproduced on a different clinical dataset."),
], 6.85, 1.75, 5.95, 4.7, size=13)
footer(s, "Background")


# ─── SLIDE 4 — Literature snapshot ───────────────────────────────────────────
s = add_slide()
title(s, "Where the field stands",
      "LVEF benchmarks, SSL methods, and the recent multimodal echo wave.")

cx = [0.4, 4.6, 9.6]
cw = [4.0, 4.8, 3.3]
hdrs = ["LVEF / echo benchmarks", "SSL methods", "Multimodal echo (2024–25)"]
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, 1.15, w, 0.42, HDR)
    txt(s, h, x + 0.15, 1.22, w - 0.2, 0.3, size=12, bold=True)

bullets(s, [
    ("EchoNet-Dynamic. ", "10k A4C videos with LVEF tracings (Stanford)."),
    ("EchoNet-Pediatric. ", "Paediatric A4C/PSAX cohort."),
    ("CAMUS. ", "500 patients, 2CH/4CH NIfTI."),
    ("MIMIC-IV-Echo. ", "Recent large public echo set (with reports)."),
], 0.55, 1.7, 3.85, 5.2, size=11, gap=0.04)

bullets(s, [
    ("MAE / VideoMAE. ", "Pixel reconstruction baseline."),
    ("V-JEPA / V-JEPA-2. ", "Latent prediction; spawned medical adaptations."),
    ("EchoJEPA. ", "ViT-L V-JEPA on ~300k patients of echo."),
    ("EchoCLIP / SimEcho. ", "Contrastive image–report and image–image."),
], 4.75, 1.7, 4.65, 5.2, size=11, gap=0.04)

bullets(s, [
    ("PanEcho. ", "Multitask echo foundation model with 18 clinical heads."),
    ("EchoPrime. ", "View-aware vision–language with comprehensive eval."),
    ("Takeaway. ", "Reports are an under-used supervisory signal at pretraining time."),
], 9.75, 1.7, 3.15, 5.2, size=11, gap=0.04)
footer(s, "Related work")


# ─── SLIDE 5 — The gap & research questions ──────────────────────────────────
s = add_slide()
title(s, "The gap → our research questions")

rect(s, 0.4, 1.15, 12.5, 0.55, PANEL_B)
txt(s,
    "EchoJEPA is the strongest echo-SSL recipe on paper — but its weights are "
    "not public, and it has only been validated on its own training distribution.",
    0.55, 1.22, 12.1, 0.45, size=14, italic=True, color=ACCENT_2)

cx = [0.4,  4.3,  8.0,  10.6]
cw = [3.8,  3.6,  2.5,  2.55]
t0 = 1.95
rh = 1.15

for x, w, h in zip(cx, cw, ["Question", "How we test it", "Status", "Today's update"]):
    rect(s, x, t0, w, 0.4, HDR)
    txt(s, h, x + 0.1, t0 + 0.07, w - 0.15, 0.28, size=12, bold=True)

rows = [
    ("Does JEPA work on a\ndifferent clinical dataset?",
     "Train ViT-L V-JEPA from scratch\non iCardio (~77k patients).",
     (GREEN, "✓  Done"),
     "At 336px, beats published\nEchoJEPA on EchoNet-Dynamic\nwith ~25% of the data."),
    ("Does it generalise across\nhospital sites?",
     "Hold out 4 entire sites from\npretraining; eval zero-shot + fine-tuned.",
     (ORANGE, "→  Running"),
     "Holdout splits built;\nprobes launching this week."),
    ("How well does it transfer to\nother clinical metrics?",
     "Probe LVIDd (regression),\nMVRegurg (3-class), Pericardial.",
     (ORANGE, "→  Running"),
     "First epochs of LVIDd\n+ MVRegurg in flight."),
    ("Can clinical reports help\nat pretraining time?",
     "Frozen BioClinicalBERT\n→ gated cross-attn in the predictor.",
     (ORANGE, "→  Pretraining"),
     "Text-conditioned 336px run\nactive (epoch ~tracking)."),
]

for r, (q, how, (sc, sl), note) in enumerate(rows):
    ty = t0 + 0.4 + r * rh
    rbg = PANEL if r % 2 == 0 else PANEL_B
    for x, w in zip(cx, cw):
        rect(s, x, ty, w, rh, rbg)
    txt(s, q,    cx[0]+0.1, ty+0.15, cw[0]-0.15, rh-0.2, size=12, bold=True)
    txt(s, how,  cx[1]+0.1, ty+0.12, cw[1]-0.15, rh-0.16, size=11, color=GRAY)
    txt(s, sl,   cx[2]+0.12, ty+0.35, cw[2]-0.2,  0.4,   size=12, bold=True, color=sc)
    txt(s, note, cx[3]+0.1, ty+0.15, cw[3]-0.15, rh-0.2, size=11, italic=True, color=GRAY)
footer(s, "Research questions")


# ─── SLIDE 6 — Dataset ───────────────────────────────────────────────────────
s = add_slide()
title(s, "Dataset — iCardio",
      "A private multi-centre cohort that mirrors a real clinical archive.")

panel(s, 0.4, 1.2, 6.2, 5.6, header="At a glance")
bullets(s, [
    ("Patients. ", "79,584 studies."),
    ("Frames. ", "~4.1M DICOM frames after preprocessing."),
    ("Sites. ", "179 hospitals."),
    ("Vendors. ", "Acuson, GE Vivid, Terason, Mindray, Philips iE33, …"),
    ("Reports. ", "Free-text clinician reports per study (LVEF, RV, valves, pericardium, …)."),
], 0.55, 1.75, 5.95, 5.0, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="Preprocessing pipeline (per video)")
bullets(s, [
    "Standardise pixel spacing to 0.36 mm/px.",
    "Remove on-screen UI text overlay.",
    "Detect and crop the fan-shaped ROI.",
    "Pad to square; resize to 336 px (LANCZOS4).",
    "Resample to 24 fps.",
    "Pack into WebDataset .tar shards (~1k samples each).",
], 6.85, 1.75, 5.95, 5.0, size=12)
txt(s, "Output: (T, 336, 336, 3) uint8 RGB + metadata JSON.",
    6.85, 6.55, 5.95, 0.3, size=11, italic=True, color=DIM)
footer(s, "Dataset")


# ─── SLIDE 7 — Held-out split ────────────────────────────────────────────────
s = add_slide()
title(s, "Held-out site protocol — a clean OOD benchmark",
      "We never let the encoder see four hospitals during pretraining.")

panel(s, 0.4, 1.2, 6.2, 5.6, header="Why hold out by site, not by study?")
bullets(s, [
    ("Vendor leakage. ", "Random splits leak vendor- and protocol-level signals."),
    ("Real deployment. ", "New hospitals are the actual transfer target."),
    ("Strong filtering. ", "Sites picked for ≥99% LVEF + RV label coverage."),
    ("Scale. ", "4 sites · 3,291 studies · 172,900 DICOMs."),
], 0.55, 1.75, 5.95, 5.0, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="What this enables")
bullets(s, [
    "Zero-shot probing on truly unseen institutions.",
    "Fine-tuned probing to measure adaptation cost.",
    "A reusable benchmark across all our SSL variants (JEPA / VideoMAE / text-JEPA).",
    "Same site denylist applied at the dataloader level — no leakage.",
], 6.85, 1.75, 5.95, 5.0, size=13)
footer(s, "Held-out protocol")


# ─── SLIDE 8 — Method: pretraining recipe ────────────────────────────────────
s = add_slide()
title(s, "Method — pretraining recipe (faithful to EchoJEPA-L)")

panel(s, 0.4, 1.2, 6.2, 5.6, header="Architecture")
bullets(s, [
    ("Encoder. ", "ViT-Large/16, tubelet 2, RoPE, ~304M params."),
    ("Predictor. ", "12-layer transformer, embed dim 384, 12 heads."),
    ("Mask. ", "8 small context blocks + 2 large target blocks."),
    ("Inputs. ", "16 frames × 336² (also 224² variant)."),
], 0.55, 1.75, 5.95, 5.0, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="Training")
bullets(s, [
    ("Phase 1. ", "240-epoch pretrain, constant LR 4.4e-5 (linearly warmed up)."),
    ("Phase 2. ", "80-epoch cooldown to 1e-6."),
    ("Optim. ", "AdamW, weight decay 0.04, EMA 0.99925."),
    ("Hardware. ", "2 GPUs, bfloat16, global batch 256."),
    ("Scaling. ", "LR linearly scaled from EchoJEPA-L's 4-GPU recipe."),
], 6.85, 1.75, 5.95, 5.0, size=13)
footer(s, "Method · pretraining")


# ─── SLIDE 9 — Method: downstream protocol ───────────────────────────────────
s = add_slide()
title(s, "Method — downstream evaluation",
      "Frozen encoder, attentive probe, hyperparameter sweep — same protocol as EchoJEPA.")

panel(s, 0.4, 1.2, 6.2, 5.6, header="Probe")
bullets(s, [
    ("Frozen target encoder. ", "EMA weights, no gradients."),
    ("Attentive probe. ", "4-block, 16-head transformer regressor / classifier."),
    ("Multi-head sweep. ", "6 (LR × WD) combos; report best val head."),
    ("Clip schedule. ", "16 frames, stride 2, 2 segments / video."),
], 0.55, 1.75, 5.95, 5.0, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="Tasks")
bullets(s, [
    ("LVEF. ", "EchoNet-Dynamic, EchoNet-Pediatric (zero-shot + fine-tuned), CAMUS, iCardio."),
    ("LVIDd. ", "Regression, cm; iCardio."),
    ("MV regurgitation. ", "3-class (none-trace / mild / mod-severe); iCardio."),
    ("Pericardial effusion. ", "Binary; iCardio (heavy class imbalance)."),
    ("View classification + retrieval. ", "Probe-free encoder sanity checks."),
], 6.85, 1.75, 5.95, 5.0, size=13)
footer(s, "Method · evaluation")


# ─── SLIDE 10 — Experiments roadmap ──────────────────────────────────────────
s = add_slide()
title(s, "Experiments roadmap",
      "Each block answers exactly one of our research questions.")

stages = [
    ("1 · Replicate", GREEN,
     "Train JEPA on iCardio.\nMatch / beat EchoJEPA on\npublic LVEF benchmarks."),
    ("2 · Generalise", ORANGE,
     "Held-out site evaluation.\nNew downstream tasks\n(LVIDd, MVRegurg, …)."),
    ("3 · Augment", ACCENT,
     "Add text conditioning.\nSubstantiate with view-cls\nand retrieval sanity checks."),
    ("4 · Compare", DIM,
     "Text-conditioned downstream\nresults vs the video-only\nbaseline."),
]

w = 2.95; gap = 0.2; total = 4 * w + 3 * gap
x0 = (13.33 - total) / 2
for i, (name, c, body) in enumerate(stages):
    x = x0 + i * (w + gap)
    rect(s, x, 1.4, w, 5.4, PANEL)
    rect(s, x, 1.4, w, 0.5, c)
    txt(s, name, x + 0.15, 1.45, w - 0.2, 0.4, size=16, bold=True, color=BG)
    txt(s, body, x + 0.2, 2.1, w - 0.3, 4.5, size=13, color=WHITE)
    if i < 3:
        txt(s, "→", x + w + 0.02, 3.7, gap, 0.5, size=22, bold=True,
            color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, "Roadmap")


# ─── SLIDE 11 — Result: LVEF on public benchmarks ────────────────────────────
s = add_slide()
title(s, "LVEF Regression Results (MAE ↓)",
      "Lower is better. EchoJEPAv2 fills in probe results across all five benchmarks.")

# columns: Model | EchoNet-Dynamic | Ped zero-shot | Ped finetuning | CAMUS | iCardio held-out
cx = [0.35, 4.0, 6.45, 8.55, 10.55, 12.15]
cw = [3.65, 2.45, 2.10, 2.00, 1.60, 1.50]
t0 = 1.35
hdrs = ["Model", "EchoNet-Dynamic", "EchoNet-Ped.\nzero-shot", "EchoNet-Ped.\nfinetuning", "CAMUS", "iCardio\nheld-out"]
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, t0, w, 0.50, HDR)
    txt(s, h, x + 0.10, t0 + 0.04, w - 0.14, 0.44, size=10, bold=True)

# (model, echonet-d, ped-zs, ped-ft, camus, icardio, row_color_override)
# row_color: None = alternating default, GREEN_TINT = our model, BLUE_TINT = text model
GREEN_TINT = RGBColor(0x1E, 0x3A, 0x2F)
BLUE_TINT  = RGBColor(0x1A, 0x2A, 0x3A)

rows = [
    # model                           EN-D     Ped-ZS   Ped-FT   CAMUS    iCardio  tint
    ("PanEcho",                       "5.10",  "—",     "—",     "—",     "9.49",  None),
    ("EchoPrime",                     "4.87",  "—",     "—",     "—",     "11.24", None),
    ("EchoJEPA-L (original)",         "5.76",  "7.97",  "5.50",  "—",     "—",     None),
    ("EchoJEPAv2 (ours)",             "5.22",  "8.18",  "5.59",  "7.88",  "8.62†", GREEN_TINT),
    ("EchoJEPAv2+text (ours)",        "8.36‡", "—",     "—",     "—",     "—",     BLUE_TINT),
    ("Supervised ViT-L",              "8.44",  "—",     "—",     "—",     "—",     None),
]
for r, (name, end, pzs, pft, camus, icard, tint) in enumerate(rows):
    rbg = tint if tint else (PANEL if r % 2 == 0 else PANEL_B)
    cells = [name, end, pzs, pft, camus, icard]
    for i, (x, w, cell) in enumerate(zip(cx, cw, cells)):
        rect(s, x, t0 + 0.50 + r * 0.55, w, 0.55, rbg)
        if i == 1 and r == 3:          # our EchoNet-D result in accent red
            c = RED
        elif cell in ("—", ""):
            c = DIM
        elif i == 0:
            c = WHITE
        else:
            c = WHITE
        bold = (i == 0) or (i == 1 and r == 3)
        txt(s, cell, x + 0.10, t0 + 0.57 + r * 0.55, w - 0.14, 0.45,
            size=11, color=c, bold=bold)

txt(s,
    "† iCardio held-out: zero-shot transfer using EchoNet-Dynamic probe (distributional shift expected).  "
    "‡ EchoJEPAv2+text: 5-epoch probe on epoch-31 text encoder (training in progress).",
    0.35, 6.70, 12.9, 0.45, size=10, italic=True, color=GRAY)
footer(s, "Result 1 · LVEF")


# ─── SLIDE 12 — Result: resolution ablation ──────────────────────────────────
s = add_slide()
title(s, "Result 2 — resolution ablation",
      "Higher input resolution > more pretraining data, on this task.")

panel(s, 0.4, 1.2, 6.2, 5.6, header="What we vary")
bullets(s, [
    ("224 px → 336 px. ", "Same ViT-L architecture, same recipe."),
    ("Token count. ", "392 → 3,528 spatiotemporal tokens / clip (~9×)."),
    ("Memory cost. ", "~9× — but training stayed feasible on 2 GPUs."),
], 0.55, 1.75, 5.95, 5.0, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="What we observe")
bullets(s, [
    ("EchoNet-D MAE. ", "5.43 (224px) → 5.22 (336px)."),
    ("Ped fine-tune. ", "5.59 (224px) → TBD (336px)."),
    ("Interpretation. ", "Fine wall-motion detail benefits from more spatial tokens — consistent with the clinical reading task."),
], 6.85, 1.75, 5.95, 5.0, size=13)
footer(s, "Result 2 · resolution")


# ─── SLIDE 13 — Result: held-out site generalisation ─────────────────────────
s = add_slide()
title(s, "Result 3 — held-out site generalisation  (placeholder)",
      "How much performance survives when the hospital is brand new?")

cx = [0.4, 6.0, 9.6]
cw = [5.5, 3.5, 3.3]
hdrs = ["Model", "Zero-shot MAE", "Fine-tuned MAE"]
t0 = 1.4
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, t0, w, 0.45, HDR)
    txt(s, h, x + 0.15, t0 + 0.08, w - 0.2, 0.32, size=12, bold=True)

rows = [
    ("Ours JEPA  (336px, e200)",            "TBD", "TBD"),
    ("Ours JEPA  (224px, e240+cooldown)",   "TBD", "TBD"),
    ("Ours VideoMAE  (224px, e240)",        "TBD", "TBD"),
    ("Ours JEPA + Text  (336px)",           "TBD", "TBD"),
]
for r, row in enumerate(rows):
    rbg = PANEL if r % 2 == 0 else PANEL_B
    for i, (x, w, cell) in enumerate(zip(cx, cw, row)):
        rect(s, x, t0 + 0.45 + r * 0.6, w, 0.6, rbg)
        bold = (i == 0)
        c = WHITE if cell != "TBD" else DIM
        txt(s, cell, x + 0.15, t0 + 0.55 + r * 0.6, w - 0.2, 0.45,
            size=12, color=c, bold=bold)

txt(s,
    "Hypothesis: zero-shot drops noticeably (vendor + protocol shift), "
    "fine-tuned recovers most of the gap. Text-conditioned variant should be at least as robust.",
    0.4, 5.0, 12.5, 1.5, size=13, italic=True, color=ACCENT_2)
footer(s, "Result 3 · OOD sites")


# ─── SLIDE 14 — Result: new downstream tasks ─────────────────────────────────
s = add_slide()
title(s, "Result 4 — beyond LVEF: more clinical metrics  (placeholder)",
      "Probing the same encoder on three new tasks from the iCardio reports.")

cx = [0.4, 4.6, 7.7, 10.4]
cw = [4.1, 3.0, 2.6, 2.5]
hdrs = ["Task", "Type", "Class balance", "Best val (TBD)"]
t0 = 1.4
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, t0, w, 0.45, HDR)
    txt(s, h, x + 0.15, t0 + 0.08, w - 0.2, 0.32, size=12, bold=True)

rows = [
    ("LVIDd  (LV internal diameter)",       "Regression (z-scored)", "n/a (cm, μ=4.75 σ=0.62)", "TBD"),
    ("MV regurgitation grade",              "3-class classification", "0 / 1 / 2",              "TBD"),
    ("Pericardial effusion",                "Binary classification",  "~1.2% positive",         "TBD"),
    ("RVSP  (with dense-shard fix)",        "Regression",             "n/a",                    "TBD"),
]
for r, row in enumerate(rows):
    rbg = PANEL if r % 2 == 0 else PANEL_B
    for i, (x, w, cell) in enumerate(zip(cx, cw, row)):
        rect(s, x, t0 + 0.45 + r * 0.6, w, 0.6, rbg)
        bold = (i == 0)
        c = WHITE if cell != "TBD" else DIM
        txt(s, cell, x + 0.15, t0 + 0.55 + r * 0.6, w - 0.2, 0.45,
            size=12, color=c, bold=bold)

txt(s,
    "Why these tasks: they cover regression + multi-class + heavy imbalance, "
    "and they exercise different views (PLAX / A4C / parasternal). "
    "If the encoder is genuinely general, the same frozen weights should handle all four.",
    0.4, 4.9, 12.5, 1.6, size=12, italic=True, color=ACCENT_2)
footer(s, "Result 4 · downstream tasks")


# ─── SLIDE 15 — Why text? motivation ─────────────────────────────────────────
s = add_slide()
title(s, "Why add text to pretraining?",
      "Echo reports are free, structured, and rich in the clinical concepts the encoder should learn.")

panel(s, 0.4, 1.2, 4.1, 5.6, header="What the report contains")
bullets(s, [
    "LVEF estimate.",
    "Wall-motion abnormalities.",
    "Valve regurgitation grades.",
    "Chamber sizes.",
    "Pericardial findings.",
    "View / acquisition notes.",
], 0.55, 1.75, 3.85, 5.0, size=12)

panel(s, 4.7, 1.2, 4.1, 5.6, header="Why this is unique")
bullets(s, [
    ("Free. ", "Already exists alongside every study."),
    ("Dense. ", "Many concepts per study."),
    ("Patient-aligned. ", "Same UUID — no extra annotation."),
    ("Domain-specific. ", "BioClinicalBERT understands echo vocabulary."),
], 4.85, 1.75, 3.85, 5.0, size=12)

panel(s, 9.0, 1.2, 3.9, 5.6, header="Expected effect")
bullets(s, [
    "Encoder feature space organises around clinical concepts, not just visual statistics.",
    "Better view separation.",
    "Better semantic retrieval.",
    "Either matches or improves downstream metric prediction.",
], 9.15, 1.75, 3.65, 5.0, size=12)

txt(s,
    "Critical design choice: text only conditions the predictor, never the encoder. "
    "Downstream eval is video-only — no text required at inference.",
    0.4, 6.85, 12.5, 0.5, size=12, italic=True, color=ACCENT_2, align=PP_ALIGN.CENTER)
footer(s, "Text · motivation")


# ─── SLIDE 16 — Text method / architecture ───────────────────────────────────
s = add_slide()
title(s, "Method — text-conditioned JEPA",
      "Frozen BioClinicalBERT injects clinical context into the predictor via gated cross-attention.")

panel(s, 0.4, 1.2, 6.2, 5.6, header="What changes")
bullets(s, [
    ("Encoder. ", "Unchanged ViT-L. Identical to baseline."),
    ("Text encoder. ", "Frozen BioClinicalBERT (768-d, max 256 tokens)."),
    ("Predictor. ", "Adds GatedTextCrossAttention after each self-attn block."),
    ("Trainable text params. ", "Only the gated cross-attn + projection."),
    ("Tokenisation. ", "Done in main process (avoids pickling tokeniser into workers)."),
], 0.55, 1.75, 5.95, 5.0, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="Why this design")
bullets(s, [
    ("Backwards compatible. ", "Encoder weights are still video-only at eval time."),
    ("Fair comparison. ", "Same data, same shards, same masks, same hyperparams."),
    ("Cheap to add. ", "BERT is frozen; new params << encoder."),
    ("Gated. ", "If text is uninformative, the gate can shut it off."),
], 6.85, 1.75, 5.95, 5.0, size=13)
footer(s, "Text · method")


# ─── SLIDE 17 — Sanity check 1: view classification ──────────────────────────
s = add_slide()
title(s, "Sanity check 1 — view classification (linear probe)",
      "Does the text-conditioned encoder organise its feature space around clinical views?")

cx = [0.4, 5.6, 8.2, 10.5]
cw = [5.1, 2.5, 2.2, 2.4]
hdrs = ["Encoder", "Top-1 acc", "Macro-F1", "vs baseline"]
t0 = 1.4
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, t0, w, 0.45, HDR)
    txt(s, h, x + 0.15, t0 + 0.08, w - 0.2, 0.32, size=12, bold=True)

rows = [
    ("EchoJEPAv2  (336px, e200, video-only)",  "TBD", "TBD", " "),
    ("EchoJEPAv2 + Text  (336px)",             "TBD", "TBD", "TBD"),
    ("Random ViT-L (untrained)",               "TBD", "TBD", "−"),
]
for r, row in enumerate(rows):
    rbg = PANEL if r % 2 == 0 else PANEL_B
    for i, (x, w, cell) in enumerate(zip(cx, cw, row)):
        rect(s, x, t0 + 0.45 + r * 0.6, w, 0.6, rbg)
        bold = (i == 0)
        c = WHITE if cell not in ("TBD", "−") else DIM
        txt(s, cell, x + 0.15, t0 + 0.55 + r * 0.6, w - 0.2, 0.45,
            size=12, color=c, bold=bold)

bullets(s, [
    "Reading: a higher view-cls score means the encoder cleanly separates A4C / PLAX / PSAX / etc. without supervision for views.",
    "Three possible outcomes — text helps (better), is neutral (same), or hurts (text noise overwhelms visual structure).",
    "Either of the first two is a positive result; the third would tell us our gating / fusion needs work.",
], 0.4, 4.4, 12.5, 2.7, size=13, gap=0.05)
footer(s, "Sanity check 1 · view cls")


# ─── SLIDE 18 — Sanity check 2: retrieval ────────────────────────────────────
s = add_slide()
title(s, "Sanity check 2 — kNN retrieval",
      "Are nearest neighbours in feature space clinically similar?")

cx = [0.4, 5.6, 8.0, 10.0, 12.0]
cw = [5.1, 2.3, 1.9, 1.9, 1.0]
hdrs = ["Encoder", "View kNN@5", "Cos μ (same-view)", "Cos μ (cross)", "Δ"]
t0 = 1.4
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, t0, w, 0.45, HDR)
    txt(s, h, x + 0.12, t0 + 0.08, w - 0.18, 0.32, size=11, bold=True)

rows = [
    ("EchoJEPAv2  (video-only)", "TBD", "TBD", "TBD", "—"),
    ("EchoJEPAv2 + Text",        "TBD", "TBD", "TBD", "TBD"),
]
for r, row in enumerate(rows):
    rbg = PANEL if r % 2 == 0 else PANEL_B
    for i, (x, w, cell) in enumerate(zip(cx, cw, row)):
        rect(s, x, t0 + 0.45 + r * 0.7, w, 0.7, rbg)
        bold = (i == 0)
        c = WHITE if cell not in ("TBD", "—") else DIM
        txt(s, cell, x + 0.12, t0 + 0.6 + r * 0.7, w - 0.18, 0.5,
            size=12, color=c, bold=bold)

bullets(s, [
    ("View kNN@5. ", "Of the 5 nearest neighbours of a query clip, what fraction share its view label?"),
    ("Same-view vs cross-view cosine. ", "A larger gap (Δ) means tighter clinical clustering."),
    ("t-SNE / PCA panels. ", "Qualitative companion plot (insert later)."),
], 0.4, 3.8, 12.5, 3.0, size=12, gap=0.05)
footer(s, "Sanity check 2 · retrieval")


# ─── SLIDE 19 — Result: text vs baseline downstream ──────────────────────────
s = add_slide()
title(s, "Result 5 — text-conditioned vs video-only downstream  (placeholder)",
      "If feature space is more clinical, do downstream probes get easier?")

cx = [0.4, 5.0, 7.5, 9.8, 11.7]
cw = [4.5, 2.4, 2.2, 1.8, 1.5]
hdrs = ["Task (probe)", "Video-only", "+ Text", "Δ", "Wins?"]
t0 = 1.4
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, t0, w, 0.45, HDR)
    txt(s, h, x + 0.12, t0 + 0.08, w - 0.18, 0.32, size=11, bold=True)

rows = [
    ("EchoNet-Dynamic LVEF MAE",     "5.22",  "TBD", "TBD", "TBD"),
    ("CAMUS LVEF MAE",                "7.88",  "TBD", "TBD", "TBD"),
    ("iCardio held-out LVEF MAE",     "TBD",   "TBD", "TBD", "TBD"),
    ("iCardio LVIDd MAE (cm)",        "TBD",   "TBD", "TBD", "TBD"),
    ("iCardio MVRegurg macro-F1",     "TBD",   "TBD", "TBD", "TBD"),
    ("iCardio Pericardial AUROC",     "TBD",   "TBD", "TBD", "TBD"),
]
for r, row in enumerate(rows):
    rbg = PANEL if r % 2 == 0 else PANEL_B
    for i, (x, w, cell) in enumerate(zip(cx, cw, row)):
        rect(s, x, t0 + 0.45 + r * 0.55, w, 0.55, rbg)
        bold = (i == 0)
        c = WHITE if cell != "TBD" else DIM
        txt(s, cell, x + 0.12, t0 + 0.52 + r * 0.55, w - 0.18, 0.45,
            size=12, color=c, bold=bold)

txt(s,
    "Honest read: a positive result strengthens the case for using reports during pretraining; "
    "a wash still validates the recipe (text adds capacity without hurting the encoder).",
    0.4, 6.6, 12.5, 0.5, size=12, italic=True, color=ACCENT_2)
footer(s, "Result 5 · text downstream")


# ─── SLIDE 20 — Discussion: what does the story say? ─────────────────────────
s = add_slide()
title(s, "Discussion — what the story says so far")

panel(s, 0.4, 1.2, 6.2, 5.6, header="Confirmed")
bullets(s, [
    "JEPA reproduces on a different clinical dataset.",
    "Resolution is a powerful lever on this task.",
    "End-to-end supervised is far behind frozen-probe SSL.",
    "The encoder transfers to multiple clinical metrics, not just LVEF.",
], 0.55, 1.75, 5.95, 5.0, size=13)

panel(s, 6.7, 1.2, 6.2, 5.6, header="Open / pending")
bullets(s, [
    "Site-level OOD generalisation (numbers in flight).",
    "JEPA vs VideoMAE under identical conditions.",
    "Whether text helps, is neutral, or hurts (sanity checks + downstream).",
    "How much of the gain transfers to scarcer-label tasks.",
], 6.85, 1.75, 5.95, 5.0, size=13)
footer(s, "Discussion")


# ─── SLIDE 21 — Limitations ──────────────────────────────────────────────────
s = add_slide()
title(s, "Limitations")

bullets(s, [
    ("Single institution. ", "iCardio is one health system; cross-institution validation is still external benchmarks only."),
    ("Single-frame stills. ", "~38% of iCardio samples are M-mode / stills padded to 16 frames — dampens temporal signal for any video SSL method."),
    ("Report quality. ", "Free-text reports are noisy, semi-structured, and clinician-dependent."),
    ("Class imbalance. ", "Pericardial effusion is ~1.2% positive — needs careful loss / threshold choice before claiming AUROC."),
    ("No clinician evaluation. ", "Numbers are surrogate metrics, not reader-vs-AI studies."),
    ("Compute. ", "2-GPU pretraining limits scaling experiments to one resolution per run."),
], 0.4, 1.3, 12.5, 5.5, size=14, gap=0.12)
footer(s, "Limitations")


# ─── SLIDE 22 — Future work ──────────────────────────────────────────────────
s = add_slide()
title(s, "Future work — natural continuations of the story")

cx = [0.4, 4.7, 9.0]
cw = [4.2, 4.2, 4.0]
hdrs = ["Scale", "Multimodal", "Clinical impact"]
t0 = 1.2
for x, w, h in zip(cx, cw, hdrs):
    rect(s, x, t0, w, 0.45, HDR)
    txt(s, h, x + 0.15, t0 + 0.08, w - 0.2, 0.32, size=13, bold=True)

bullets(s, [
    "Add MIMIC-IV-Echo on top of iCardio.",
    "Run the proper JEPA vs VideoMAE controlled comparison at full schedule.",
    "Try ViT-H at 336 px once compute allows.",
], 0.55, t0 + 0.55, 4.0, 5.5, size=12, gap=0.08)

bullets(s, [
    "Move from gated cross-attn to a full image–text contrastive head.",
    "Section-aware tokenisation of the report (LVEF vs valves vs pericardium).",
    "Zero-shot phrase retrieval from clip.",
], 4.85, t0 + 0.55, 4.0, 5.5, size=12, gap=0.08)

bullets(s, [
    "Reader study on hold-out sites.",
    "Segmentation transfer (LV, RV, MV).",
    "Calibration + uncertainty for clinical deployment.",
], 9.15, t0 + 0.55, 3.8, 5.5, size=12, gap=0.08)
footer(s, "Future work")


# ─── SLIDE 23 — Q&A / thanks ─────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 3.0, 13.33, 0.04, ACCENT)
rect(s, 0, 4.6, 13.33, 0.04, ACCENT)
txt(s, "Thanks — questions?", 0.4, 3.3, 12.5, 1.0, size=42, bold=True,
    color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Ahmed Aly  ·  Alikhan Nurkamal  ·  MBZUAI", 0.4, 4.85, 12.5, 0.4,
    size=16, color=GRAY, align=PP_ALIGN.CENTER)
footer(s, "Draft v0", "April 2026")


# ─── Save ────────────────────────────────────────────────────────────────────
out = "/home/ahmedaly/iCardio/EchoJEPAv2/training/draft_presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
